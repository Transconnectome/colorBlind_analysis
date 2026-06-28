#!/usr/bin/env python3
"""Build ICML SD4H poster — 24 W x 36 H portrait (conference hard cap).
Spine = method-as-template (structured-distortion inference -> analytic inversion).
Results band = 3 panels: A geometry exists -> B invertible filter -> C 2nd-MRI validation.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

_HERE = os.path.dirname(os.path.abspath(__file__))
ASSET = os.path.join(_HERE, "assets")
FIG   = os.path.join(_HERE, "fig")
OUT   = "/Users/jinilkim/Library/CloudStorage/OneDrive-Personal/Projects/colorBlind_analysis/docs/ICML_workshop/poster_ICML_SD4H.pptx"

# ---------- palette ----------
NAVY   = RGBColor(0x13, 0x29, 0x4B)
BLUE   = RGBColor(0x2E, 0x5A, 0x9C)
INK    = RGBColor(0x22, 0x22, 0x22)
GRAY   = RGBColor(0x55, 0x55, 0x55)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xE8, 0x76, 0x2C)   # deutan sub-08
SBLUE  = RGBColor(0x2E, 0x6F, 0xBF)   # protan sub-09
GREENF = RGBColor(0x3A, 0x8E, 0x4A)
LIGHT  = RGBColor(0xF3, 0xF5, 0xF8)
PANEL  = RGBColor(0xF7, 0xF9, 0xFB)
TABHEAD= RGBColor(0xE3, 0xE9, 0xF2)
CELLLN = RGBColor(0xCC, 0xD3, 0xDD)
FONT = "Arial"

# ---------- geometry (inches) ----------
PW, PH = 24.0, 36.0
M  = 0.55
G  = 0.40
CW = (PW - 2*M - G)/2          # = 11.25
LX = M
RX = M + CW + G

prs = Presentation()
prs.slide_width  = Inches(PW)
prs.slide_height = Inches(PH)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# ---------- helpers ----------
def rect(x, y, w, h, color=None, line=None, line_w=None, round=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if round:
        try: shp.adjustments[0] = 0.05
        except Exception: pass
    if color is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = color
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(line_w or 1)
    shp.shadow.inherit = False
    return shp

def textbox(x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            wrap=True, fill=None, line=None, line_w=1.0, pad=0.10):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    for m in ('margin_left','margin_right','margin_top','margin_bottom'):
        setattr(tf, m, Inches(pad))
    if fill is not None: tb.fill.solid(); tb.fill.fore_color.rgb = fill
    else: tb.fill.background()
    if line is not None: tb.line.color.rgb = line; tb.line.width = Pt(line_w)
    else: tb.line.fill.background()
    tb.shadow.inherit = False
    first = True
    for p in paras:
        txt=p.get('t',''); size=p.get('s',15); bold=p.get('b',False)
        color=p.get('c',INK); lvl=p.get('lvl',0); sa=p.get('sa',4); sb=p.get('sb',0)
        bullet=p.get('bul',None); align_p=p.get('al',align); lh=p.get('lh',1.0)
        runs=p.get('runs',None)
        para = tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        para.alignment=align_p; para.space_after=Pt(sa); para.space_before=Pt(sb)
        try: para.line_spacing=lh
        except Exception: pass
        indent=0.0
        if bullet is not None: prefix=bullet+'  '; indent=0.26*(lvl+1)
        else: prefix=''
        if runs:
            if prefix:
                r0=para.add_run(); r0.text=prefix
                r0.font.size=Pt(size); r0.font.bold=bold; r0.font.name=FONT; r0.font.color.rgb=color
            for rt,rb,rc in runs:
                r=para.add_run(); r.text=rt
                r.font.size=Pt(size); r.font.bold=rb; r.font.name=FONT
                r.font.color.rgb = rc if rc is not None else color
        else:
            r=para.add_run(); r.text=prefix+txt
            r.font.size=Pt(size); r.font.bold=bold; r.font.name=FONT; r.font.color.rgb=color
        pPr=para._p.get_or_add_pPr()
        pPr.set('marL', str(int(indent*914400)))
        pPr.set('indent', str(-int(0.26*914400)) if bullet is not None else '0')
    return tb

def bar(x, y, w, title, h=0.62, size=19, fill=NAVY):
    b=rect(x,y,w,h,fill)
    tf=b.text_frame; tf.vertical_anchor=MSO_ANCHOR.MIDDLE; tf.word_wrap=True
    tf.margin_top=Inches(0.01); tf.margin_bottom=Inches(0.01)
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=title.upper(); r.font.size=Pt(size); r.font.bold=True
    r.font.name=FONT; r.font.color.rgb=WHITE
    return y+h

def section(x, y, w, title, body_paras, bar_h=0.62, body_h=2.0,
            body_fill=None, title_size=19):
    bar(x,y,w,title,h=bar_h,size=title_size)
    if body_paras is not None:
        textbox(x, y+bar_h, w, body_h, body_paras, fill=body_fill,
                line=CELLLN, line_w=0.75)
    return y + bar_h + (body_h if body_paras is not None else 0)

def pic_fit(path, x, y, w, h=None, align='center'):
    im=Image.open(path); iw,ih=im.size; ar=ih/iw
    dw=w; dh=w*ar
    if h is not None and dh>h: dh=h; dw=h/ar
    px = x+(w-dw)/2 if align=='center' else x
    slide.shapes.add_picture(path, Inches(px), Inches(y), Inches(dw), Inches(dh))
    return y+dh, px, dw

def table(x, y, w, rows, colw, rh=0.5, fs=13, head_navy=True, firstcol_colors=None):
    """rows: list of tuples; first row is header if head_navy."""
    for i,row in enumerate(rows):
        yr=y+i*rh; xx=x; hd=(i==0 and head_navy)
        for j,val in enumerate(row):
            fillc = TABHEAD if hd else (LIGHT if i%2 else WHITE)
            cell=rect(xx,yr,colw[j],rh, fillc, line=CELLLN, line_w=0.75)
            ctf=cell.text_frame; ctf.vertical_anchor=MSO_ANCHOR.MIDDLE; ctf.word_wrap=True
            ctf.margin_left=Inches(0.06); ctf.margin_right=Inches(0.04)
            ctf.margin_top=Inches(0); ctf.margin_bottom=Inches(0)
            cp=ctf.paragraphs[0]; cp.alignment=PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER
            cr=cp.add_run(); cr.text=val; cr.font.size=Pt(fs)
            cr.font.bold = hd or (j==0 and not hd)
            cr.font.name=FONT
            if j==0 and firstcol_colors and not hd:
                cr.font.color.rgb = firstcol_colors[i]
            else:
                cr.font.color.rgb = NAVY if hd else INK
            xx+=colw[j]
    return y+len(rows)*rh

# =====================================================================
# HEADER
# =====================================================================
TOP=0.35
crest=f"{ASSET}/snu_crest.png"
ch=2.05
slide.shapes.add_picture(crest, Inches(M+0.05), Inches(TOP+0.05), height=Inches(ch))
qr_w=1.35; qy=TOP+0.05
qx2=PW-M-qr_w; qx1=qx2-qr_w-0.3
slide.shapes.add_picture(f"{ASSET}/qr_code.png",  Inches(qx1), Inches(qy), Inches(qr_w), Inches(qr_w))
slide.shapes.add_picture(f"{ASSET}/qr_repro.png", Inches(qx2), Inches(qy), Inches(qr_w), Inches(qr_w))
textbox(qx1, qy+qr_w-0.04, qr_w, 0.34, [{'t':'Code & data','s':10.5,'b':True,'c':NAVY,'al':PP_ALIGN.CENTER,'sa':0}])
textbox(qx2, qy+qr_w-0.04, qr_w, 0.34, [{'t':'Repro notebooks','s':10.5,'b':True,'c':NAVY,'al':PP_ALIGN.CENTER,'sa':0}])

tx=M+ch+0.35; tw=qx1-tx-0.25
textbox(tx, TOP-0.05, tw, 1.45, [
    {'t':'Inferring Individualized Color-Vision Distortions','s':33,'b':True,'c':NAVY,'al':PP_ALIGN.CENTER,'sa':2,'lh':1.0},
    {'t':'from fMRI Hue-Representation Geometry','s':33,'b':True,'c':NAVY,'al':PP_ALIGN.CENTER,'sa':0,'lh':1.0},
], anchor=MSO_ANCHOR.TOP)
textbox(tx, TOP+1.42, tw, 0.42, [
    {'runs':[('Jinil Kim',True,INK),('¹  ',False,GRAY),('Albert Minkue Cho',True,INK),('²  ',False,GRAY),
             ('Jungwoo Seo',True,INK),('³  ',False,GRAY),('Jiook Cha',True,INK),('² ⁵ ⁶ *',False,GRAY)],
     's':16,'al':PP_ALIGN.CENTER,'sa':0}], anchor=MSO_ANCHOR.MIDDLE)
textbox(tx, TOP+1.84, tw, 0.36, [
    {'t':'¹Linguistics · ²Psychology · ³Computer Science & Engineering · ⁵Brain & Cognitive Sciences · '
         '⁶Interdisciplinary Program in AI — Seoul National University',
     's':11,'c':GRAY,'al':PP_ALIGN.CENTER,'sa':0,'lh':1.0}], anchor=MSO_ANCHOR.TOP)
rib_y=TOP+2.22
rb=rect(tx, rib_y, tw, 0.44, BLUE, round=True)
rtf=rb.text_frame; rtf.vertical_anchor=MSO_ANCHOR.MIDDLE
rp=rtf.paragraphs[0]; rp.alignment=PP_ALIGN.CENTER
rr=rp.add_run(); rr.text='Structured Data for Health (SD4H) Workshop  ·  ICML 2026'
rr.font.size=Pt(14.5); rr.font.bold=True; rr.font.name=FONT; rr.font.color.rgb=WHITE

HEAD_BOT=TOP+2.78
rect(M, HEAD_BOT, PW-2*M, 0.045, NAVY)

# =====================================================================
# TAKEAWAY strip (method-as-template spine)
# =====================================================================
tk_y=HEAD_BOT+0.16
tkb=rect(M, tk_y, PW-2*M, 0.92, NAVY, round=True)
textbox(M+0.25, tk_y, PW-2*M-0.5, 0.92, [
    {'runs':[('Many health conditions leave ',False,WHITE),('structured — not absent — signatures',True,RGBColor(0xFF,0xC8,0x8A)),
             (' in high-dimensional measurements. We turn one into a ',False,WHITE),
             ('low-dimensional, invertible, individualized correction',True,RGBColor(0xFF,0xC8,0x8A)),
             (' — read off a person\'s own cortex. Testbed: color-vision deficiency (fMRI).',False,WHITE)],
     's':16.5,'al':PP_ALIGN.CENTER,'sa':0,'lh':1.02}], anchor=MSO_ANCHOR.MIDDLE)

# =====================================================================
# BODY (2-col)
# =====================================================================
y0=tk_y+0.92+0.22
GAP=0.22

# ---------- RIGHT COLUMN: Framework fig + Models table ----------
ry=y0
ry=bar(RX, ry, CW, 'Framework — diagnose → fit → invert')
fy,_,_=pic_fit(f"{FIG}/fig1_pipeline_revise.png", RX+0.1, ry+0.12, CW-0.2, h=6.0)
textbox(RX, fy+0.04, CW, 0.78, [
    {'runs':[('Top:',True,INK),(' CVD distorts the continuous hue circle — neighbors collapse or expand — though classification stays intact. ',False,GRAY),
             ('Bottom:',True,INK),(' per-subject pipeline — diagnose, fit a 2-parameter cortical model to behavioral + neural loss, invert to a stimulus-space filter.',False,GRAY)],
     's':11.5,'sa':0,'lh':1.0}])
ry=fy+0.84+GAP

# MODELS & FITTING
ry=bar(RX, ry, CW, 'Models & fitting')
mrows=[('Mechanism','DOF','Form'),
       ('Retinal cone-shift (Machado)','1','Δλ'),
       ('Retinal + cortical (R+C)','1','δθ = (2−g)·δθ_Machado'),
       ('Cortical 2-Component','2','θ′ = θ + β_s cos(θ−90°) + β_c cos(θ−θ_conf)')]
mcolw=[CW*0.40, CW*0.10, CW*0.50]
mty=table(RX, ry+0.1, CW, mrows, mcolw, rh=0.58, fs=12.5,
          firstcol_colors=[NAVY,INK,INK,NAVY])
# bold the 2-Component row name -> handled by table firstcol; emphasize via color NAVY
textbox(RX, mty+0.08, CW, 2.0, [
    {'runs':[('β_s',True,INK),(' = S-cone axis, ',False,INK),('β_c',True,INK),(' = confusion axis (θ_conf 16° protan / 150° deutan).',False,INK)],'s':13,'bul':'•','sa':5,'lh':1.0},
    {'runs':[('Loss = ',False,INK),('behavioral JND-ratio rank',True,BLUE),(' + ',False,INK),('neural ΔRDM cosine',True,BLUE),
             (', selected by ',False,INK),('held-out generalization',True,INK),(' (7-fold leave-one-HC-out) — not in-sample significance.',False,INK)],'s':13,'bul':'•','sa':5,'lh':1.0},
    {'runs':[('Invert:',True,BLUE),(' the 2-Component map is bijective → per-hue pre-image θ̃_k by root-finding; correction δθ_k = θ̃_k − θ_k.',False,INK)],'s':13,'bul':'•','sa':0,'lh':1.0},
])
RIGHT_BOT=mty+0.1+2.0

# ---------- LEFT COLUMN: 4 text sections ----------
_bodies=[1.75, 2.85, 2.55, 3.25]   # Key, Background, ThisWork, Methods
_fixed=4*0.62 + sum(_bodies)
GAP_L=max(GAP, (RIGHT_BOT - y0 - _fixed)/3.0)
ly=y0
ly=section(LX, ly, CW, 'Key Claim', [
    {'runs':[('In color-vision deficiency (CVD), the ',False,INK),('continuous geometry',True,BLUE),
             (' of cortical hue representations is distorted while ',False,INK),
             ('categorical recognition is preserved.',True,BLUE)],'s':16,'sa':6,'lh':1.04},
    {'runs':[('We quantify each individual\'s distortion, fit a ',False,INK),('2-parameter cortical model',True,INK),
             (', and ',False,INK),('analytically invert',True,BLUE),(' it into a per-person correction filter — exact for every displayed hue (',False,INK),
             ('residual < 0.001°',True,INK),(').',False,INK)],'s':16,'sa':0,'lh':1.04},
], body_h=_bodies[0], body_fill=LIGHT)
ly+=GAP_L
ly=section(LX, ly, CW, 'Background & Gap', [
    {'runs':[('Many health phenotypes are ',False,INK),('structured distortions',True,INK),
             (' of a high-dimensional response space (cardiac shape, gene-expression latent shifts). CVD is a tractable, ground-truth-checkable instance.',False,INK)],'s':14.5,'bul':'•','sa':5,'lh':1.0},
    {'t':'Such structured signatures are rarely turned into individualized corrections. The closest analogue (cochlear WHIS) simulates loss — it does not correct it.','s':14.5,'bul':'•','sa':5,'lh':1.0},
    {'runs':[('CVD ≈ 8% of males; cortical distortion ',False,INK),('varies across individuals',True,INK),
             (' even within one diagnostic category. Generic notch filters shift appearance but barely change discrimination.',False,INK)],'s':14.5,'bul':'•','sa':5,'lh':1.0},
    {'runs':[('Gap →',True,BLUE),(' an individualized, invertible correction read from each person\'s own cortex.',False,INK)],'s':14.5,'bul':'•','sa':0,'lh':1.0},
], body_h=_bodies[1])
ly+=GAP_L
ly=section(LX, ly, CW, 'This Work — the template, in 3 steps', [
    {'runs':[('① Represent  ',True,BLUE),('cast the distorted neural geometry as a structured-distortion object — an ',False,INK),
             ('interpolation-vulnerability profile',True,INK),(' (v) + ',False,INK),('pairwise ΔRDM',True,INK),(' — invisible to classification.',False,INK)],'s':14.5,'sa':6,'lh':1.0},
    {'runs':[('② Fit  ',True,BLUE),('a low-dimensional ',False,INK),('2-DOF interpretable',True,INK),
             (' model, fusing fragmented behavioral + neural views, selected by ',False,INK),('held-out loss',True,INK),('.',False,INK)],'s':14.5,'sa':6,'lh':1.0},
    {'runs':[('③ Invert  ',True,BLUE),('analytically invert it into a ',False,INK),('deployable, subject-specific',True,INK),(' correction — and test whether it works (2nd MRI).',False,INK)],'s':14.5,'sa':0,'lh':1.0},
], body_h=_bodies[2], body_fill=LIGHT)
ly+=GAP_L
ly=section(LX, ly, CW, 'Methods — data & readouts', [
    {'runs':[('N = 10',True,INK),(': 7 HC (3F, 22.7±2.5 y) + 3 CVD (2 deutan, 1 protan), Ishihara-confirmed. Small sample → ',False,INK),
             ('single-case statistics',True,INK),(' (Crawford–Howell).',False,INK)],'s':14,'bul':'•','sa':5,'lh':1.0},
    {'runs':[('Stimuli:',True,INK),(' 8 isoluminant hues (CIE L*a*b*, L*=75, chroma 40), RSVP, 6 runs, Siemens 3T → fMRIPrep → V1/V2/V3/hV4 (Wang 50%); FIR-GLM, Procrustes-aligned. JND outside scanner.',False,INK)],'s':14,'bul':'•','sa':5,'lh':1.0},
    {'runs':[('One forward-encoding model',True,INK),(' (F = 6 channels), two hold-outs: ',False,INK),
             ('LORO',True,BLUE),(' (hold a run → classification) vs ',False,INK),('LOCO',True,BLUE),(' (hold a hue → interpolation).',False,INK)],'s':14,'bul':'•','sa':5,'lh':1.0},
    {'runs':[('Signatures:',True,INK),(' vulnerability v ∈ [0,1]⁸ at hV4 + ΔRDM = RDM(CVD) − mean RDM(HC), 28 pairs.',False,INK)],'s':14,'bul':'•','sa':0,'lh':1.0},
], body_h=_bodies[3])
LEFT_BOT=ly

# =====================================================================
# RESULTS band (full width, 3 panels)
# =====================================================================
RES_TOP=max(LEFT_BOT, RIGHT_BOT)+0.28
ry=bar(M, RES_TOP, PW-2*M, 'Results', h=0.78, size=24)
ry+=0.18

pw3=(PW-2*M-2*G)/3.0
px=[M, M+pw3+G, M+2*(pw3+G)]
PANEL_H=10.6
for xx in px:
    rect(xx, ry, pw3, PANEL_H, PANEL, line=RGBColor(0xD7,0xDE,0xE7), line_w=1.0)

def subbar(x,w,text,y,col=BLUE):
    sb=rect(x,y,w,0.6,col)
    t=sb.text_frame; t.vertical_anchor=MSO_ANCHOR.MIDDLE; t.word_wrap=True
    pp=t.paragraphs[0]; pp.alignment=PP_ALIGN.CENTER
    rr=pp.add_run(); rr.text=text; rr.font.size=Pt(15.5); rr.font.bold=True; rr.font.name=FONT; rr.font.color.rgb=WHITE
    return y+0.6

pad=0.16
# ---- Panel A: structured distortion exists ----
ax=px[0]; aw=pw3
ay=subbar(ax, aw, 'A.  A structured cortical distortion', ry)
# SRM hue-geometry wheels (OHBM explainer) with colored subject headers
third=(aw-2*pad)/3.0
for k,(lab,col) in enumerate([('HC  (n = 7)',GREENF),('Sub-08 deutan · V2',ORANGE),('Sub-09 protan · V1',SBLUE)]):
    textbox(ax+pad+k*third, ay+0.1, third, 0.34,
            [{'t':lab,'s':11.5,'b':True,'c':col,'al':PP_ALIGN.CENTER,'sa':0,'lh':0.95}], anchor=MSO_ANCHOR.MIDDLE)
awy,_,_=pic_fit(f"{FIG}/srm_wheels.png", ax+pad, ay+0.46, aw-2*pad)
textbox(ax+pad, awy+0.0, aw-2*pad, 0.5, [
    {'runs':[('SRM-aligned hue geometry.',True,GRAY),(' Bold dots = each CVD, displaced from the HC mean; deviation peaks at a distinct early ROI per subject.',False,GRAY)],'s':10.5,'sa':0,'lh':0.95}])
# disparity-by-ROI (portrait) centered + bullets
ady,adx,adw=pic_fit(f"{FIG}/disparity_roi.png", ax+pad, awy+0.58, aw-2*pad, h=4.5)
textbox(ax+pad, ady+0.0, aw-2*pad, 0.4, [
    {'t':'SRM disparity by ROI — each CVD exceeds the HC ±1 SD band at one ROI.','s':10.5,'c':GRAY,'al':PP_ALIGN.CENTER,'sa':0,'lh':0.95}])
textbox(ax+pad, ady+0.46, aw-2*pad, 3.4, [
    {'runs':[('Discrimination preserved.',True,BLUE),(' Cross-run (LORO) classification: no HC–CVD difference (pooled p=0.668).',False,INK)],'s':13,'bul':'•','sa':7,'lh':1.0},
    {'runs':[('Interpolation impaired.',True,BLUE),(' Only hV4 interpolates above chance in HC (adjacent acc 0.47, p=0.044); both CVD below — deficit on S-cone hues.',False,INK)],'s':13,'bul':'•','sa':7,'lh':1.0},
    {'runs':[('Distinct ROI per subject.',True,BLUE),(' Disparity elevated at Sub-08 ',False,INK),('V2 (p=0.040)',True,ORANGE),(', Sub-09 ',False,INK),('V1 (p=0.007)',True,SBLUE),(' — Crawford–Howell.',False,INK)],'s':13,'bul':'•','sa':0,'lh':1.0},
])

# ---- Panel B: individualized invertible filter ----
bx=px[1]; bw=pw3
by=subbar(bx, bw, 'B.  An individualized, invertible filter', ry)
bfy,_,_=pic_fit(f"{FIG}/fig2_landscape_filter-1.png", bx+pad, by+0.2, bw-2*pad, h=4.95)
textbox(bx+pad, bfy+0.0, bw-2*pad, 0.5, [
    {'runs':[('Per-subject 2-Component fit (left: loss landscape over (β_s, β_c); right: inverted filter, original vs corrected hues).',False,GRAY)],'s':10.5,'sa':0,'lh':0.95}])
# Table 1
tby=bfy+0.6
textbox(bx+pad, tby, bw-2*pad, 0.34, [{'t':'Table 1 — selected 2-Component fit (hV4)','s':12.5,'b':True,'c':NAVY,'sa':0}])
tby+=0.38
t1=[('Subject','(β_s, β_c)','Test loss','|δθ|'),
    ('Sub-08 deutan','(+6°, −42°)','−2.36','26.3°'),
    ('Sub-09 protan','(+2°, +24°)','−1.54','16.2°')]
t1w=[ (bw-2*pad)*c for c in (0.34,0.26,0.22,0.18)]
tby=table(bx+pad, tby, bw-2*pad, t1, t1w, rh=0.52, fs=12.5,
          firstcol_colors=[NAVY,ORANGE,SBLUE])
textbox(bx+pad, tby+0.18, bw-2*pad, 3.6, [
    {'runs':[('Structure beats cone-shift.',True,BLUE),(' R+C overcompensates (gain pinned at ceiling) and is ',False,INK),('non-invertible at protan severity',True,INK),(' (4/8 hues lack a pre-image). 2-Component fits the grid interior for both.',False,INK)],'s':13.5,'bul':'•','sa':8,'lh':1.0},
    {'runs':[('Neural term carries what behavior cannot.',True,BLUE),(' For Sub-09, behavioral-only beats the null in 3/7 folds; adding neural ΔRDM → 7/7 + a stable confusion-axis rotation.',False,INK)],'s':13.5,'bul':'•','sa':8,'lh':1.0},
    {'runs':[('Exact inversion.',True,BLUE),(' A pre-image for all 8 hues (residual < 0.001°) → a directly implementable 8-point correction LUT; filters differ by subtype (β_c −42° deutan vs +24° protan).',False,INK)],'s':13.5,'bul':'•','sa':0,'lh':1.0},
])

# ---- Panel C: does it work? (2nd MRI) ----
cx=px[2]; cw=pw3
cy=subbar(cx, cw, 'C.  Does the filter work?  (2nd MRI)', ry, col=ORANGE)
cfy,_,_=pic_fit(f"{FIG}/fig_validation.png", cx+pad, cy+0.18, cw-2*pad)
textbox(cx+pad, cfy+0.02, cw-2*pad, 0.5, [
    {'runs':[('CVD subject views the 8 hues through each filter in-scanner (4 runs/condition); ',False,GRAY),('Optimal',True,ORANGE),(' = personalized inverse filter, ',False,GRAY),('macOS',True,INK),(' = deployed generic filter.',False,GRAY)],'s':10.5,'sa':0,'lh':0.95}])
ccy=cfy+0.62
textbox(cx+pad, ccy, cw-2*pad, 5.2, [
    {'runs':[('Behavioral parity.',True,GREENF),(' The personalized filter restores discrimination to HC level (mean JND 0.19→0.08; HC 0.10), ',False,INK),
             ('on par with the deployed macOS filter',True,INK),(' (Wilcoxon p=0.84).',False,INK)],'s':13.5,'bul':'•','sa':8,'lh':1.0},
    {'runs':[('Neural superiority.',True,ORANGE),(' Yet ',False,INK),('only the model-derived filter',True,INK),(' restores hV4 interpolation geometry toward HC (ρ +0.18 ≈ HC 0.21); the macOS filter pushes it further away (−0.39). Δ = 3.2 HC-SD, run-separated.',False,INK)],'s':13.5,'bul':'•','sa':8,'lh':1.0},
    {'runs':[('Scope.',True,GRAY),(' Descriptive proof-of-concept (Sub-08; n=4/condition, no permutation; macOS vs PsychoPy rendering differs). ',False,GRAY),('Sub-09 added 2026-06-29 → N=2.',True,SBLUE)],'s':13,'bul':'•','sa':0,'lh':1.0},
])

RES_BOT=ry+PANEL_H

# =====================================================================
# FOOTER: Conclusion + References
# =====================================================================
BOT_TOP=RES_BOT+0.3
section(LX, BOT_TOP, CW, 'Conclusion & next step', [
    {'runs':[('A worked template:',True,BLUE),(' quantify a structured fMRI distortion, fit a parsimonious interpretable model, ',False,INK),('analytically invert',True,BLUE),(' it into a deployable per-person filter.',False,INK)],'s':14,'bul':'•','sa':6,'lh':1.0},
    {'runs':[('Validated (proof-of-concept):',True,ORANGE),(' the personalized filter matches a deployed generic filter behaviorally, and ',False,INK),('uniquely repairs the cortical hue geometry',True,INK),(' (hV4). Sub-09 in progress.',False,INK)],'s':14,'bul':'•','sa':6,'lh':1.0},
    {'runs':[('For SD4H:',True,BLUE),(' mapping structured signatures onto low-dimensional, invertible features turns descriptive phenotypes into individualized intervention targets.',False,INK)],'s':14,'bul':'•','sa':0,'lh':1.0},
], body_h=3.5, body_fill=LIGHT)
section(RX, BOT_TOP, CW, 'References', [
    {'t':'1.  Brouwer & Heeger. Decoding and reconstructing color from human visual cortex. J Neurosci 29, 13992 (2009).','s':12.5,'sa':4,'lh':1.0},
    {'t':'2.  Machado, Oliveira & Fernandes. A physiologically-based model for simulation of CVD. IEEE TVCG 15, 1291 (2009).','s':12.5,'sa':4,'lh':1.0},
    {'t':'3.  Tregillus et al. Color compensation in anomalous trichromats assessed with fMRI. Curr Biol 31, 936 (2021).','s':12.5,'sa':4,'lh':1.0},
    {'t':'4.  Kriegeskorte, Mur & Bandettini. Representational similarity analysis. Front Syst Neurosci 2:4 (2008).','s':12.5,'sa':4,'lh':1.0},
    {'t':'5.  Irino. Hearing impairment simulator (WHIS). IEEE Access 11, 78419 (2023).','s':12.5,'sa':4,'lh':1.0},
    {'t':'6.  Crawford & Howell. Comparing an individual\'s score against norms from small samples. Clin Neuropsychol 12, 482 (1998).','s':12.5,'sa':4,'lh':1.0},
    {'runs':[('Correspondence:',True,NAVY),(' connectome@snu.ac.kr   ·   github.com/Transconnectome/colorBlind_analysis',False,GRAY)],'s':12,'sa':0,'sb':5,'lh':1.0},
], body_h=3.5)

prs.save(OUT)
print("saved", OUT)
print("RIGHT_BOT=%.2f LEFT_BOT=%.2f RES_TOP=%.2f RES_BOT=%.2f BOT_TOP=%.2f PH=%.2f"%(
    RIGHT_BOT,LEFT_BOT,RES_TOP,RES_BOT,BOT_TOP,PH))
print("FOOTER_BOT=%.2f (cap %.2f)"%(BOT_TOP+0.62+4.4, PH))
