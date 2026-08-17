#!/usr/bin/env python
"""Revise meeting_2026-06-01.pptx with s18 held-out test-loss results.

Edits (user-approved 2026-06-02):
 - slide4 (M2.5 criteria): add held-out test-loss as selection criterion (e).
 - NEW slide (after slide6): held-out test-loss = is the stable value GOOD.
 - slide7 (RQ4): keep S09 headline + add honest caveat; note S08 triangulation.
 - slide8 (RQ3): add broad-basin reconciliation line.
 - slide12 (limits): split generalization into HC-pool (done) vs CVD (Phase3).
 - slide6: subtitle bridge.
 - renumber page labels after insertion.

Run in srm env. Backup already at meeting_2026-06-01.BACKUP.pptx.
"""
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

SRC = "meeting_2026-06-01.pptx"
prs = Presentation(SRC)

DARK = "1A1A1A"; GREEN = "0E6B4F"; AMBER = "D47E00"; WHITE = "FFFFFF"
BLUE_SUB = "A8CCEC"; GRAY = "7F8C8D"


def shp(slide, sid):
    return [x for x in slide.shapes if x.shape_id == sid][0]


def set_run_text(p_el, text):
    """Set first <a:r>'s text, drop extra runs (preserve run formatting)."""
    rs = p_el.findall(qn('a:r'))
    if not rs:
        return
    rs[0].find(qn('a:t')).text = text
    for r in rs[1:]:
        p_el.remove(r)


def append_cloned(shape, specs):
    """specs: list of (template_para_idx, text). Clone each template paragraph,
    set text, append in order at end of the text frame (preserves formatting)."""
    tf = shape.text_frame
    anchor = tf.paragraphs[-1]._p
    for tmpl_idx, text in specs:
        new = deepcopy(tf.paragraphs[tmpl_idx]._p)
        set_run_text(new, text)
        anchor.addnext(new)
        anchor = new


def rebuild_textbox(shape, specs):
    """Rebuild a text box from scratch. specs: (text,size,bold,colorhex)."""
    tf = shape.text_frame
    tf.clear()
    for i, (text, size, bold, col) in enumerate(specs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run(); run.text = text
        run.font.size = Pt(size); run.font.bold = bold
        run.font.color.rgb = RGBColor.from_string(col)


def set_cell(cell, text, size=11, bold=False, col=DARK):
    cell.text = text
    para = cell.text_frame.paragraphs[0]
    for run in para.runs:
        run.font.size = Pt(size); run.font.bold = bold
        run.font.color.rgb = RGBColor.from_string(col)


def find_slide_by_title(substr):
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame and substr in sh.text_frame.text:
                return s
    return None


# ---------------------------------------------------------------- slide 6 bridge
s6 = prs.slides[5]
sub6 = shp(s6, 5)
set_run_text(sub6.text_frame.paragraphs[1]._p if len(sub6.text_frame.paragraphs) > 1
             else sub6.text_frame.paragraphs[0]._p,
             "v6 PCA 45° RDM  |  N=300 HC resamples + strict 7-fold LOO  →  "
             "stability = reproducibility (goodness: next slide)")

# ---------------------------------------------------------------- slide 4 criteria
s4 = prs.slides[3]
crit = shp(s4, 11)
# add after p8 ('Strict LOO range ↓'); clone p4 (bold header) + p5 (normal body)
# but append_cloned adds at end; we want after p8. Insert manually after p8.
tf = crit.text_frame
p8 = tf.paragraphs[8]._p
hdr = deepcopy(tf.paragraphs[4]._p); set_run_text(hdr, "Generalization: Held-out test-loss")
bod = deepcopy(tf.paragraphs[5]._p); set_run_text(bod, "ΔL vs no-corr.(0,0) < 0, held-out HC (s18) = criterion (e)")
p8.addnext(hdr); hdr.addnext(bod)

# ---------------------------------------------------------------- NEW slide (clone s6)
def clone_after(src_slide, dest_index):
    """Deep-copy src_slide shapes (skip pictures) onto a new blank slide,
    then move it to dest_index (0-based) in the slide order."""
    layout = src_slide.slide_layout
    new = prs.slides.add_slide(layout)
    for ph in list(new.shapes):           # remove layout placeholders
        ph._element.getparent().remove(ph._element)
    for sh in src_slide.shapes:
        if sh.shape_type == 13:            # PICTURE -> skip (broken rel)
            continue
        new.shapes._spTree.append(deepcopy(sh._element))
    # move new slide (currently last) to dest_index
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    sldIdLst.remove(ids[-1])
    sldIdLst.insert(dest_index, ids[-1])
    return new

ns = clone_after(s6, 6)   # insert as 7th slide (index 6)

# title / number / subtitle
set_run_text(shp(ns, 3).text_frame.paragraphs[1]._p,
             "Held-out Test-loss — 안정적 값은 '좋은' 값인가  (기준 e)")
set_run_text(shp(ns, 4).text_frame.paragraphs[1]._p, "7")
set_run_text(shp(ns, 5).text_frame.paragraphs[1]._p,
             "s18 · leave-one-HC-out 7-fold  |  per-fold REFIT (not plug-in)  |  "
             "ΔL vs no-correction (0,0)")

# table[6]: repopulate 4x7 as ΔL table, widen to full width
tbsh = shp(ns, 6)
tbsh.left = Inches(0.40); tbsh.top = Inches(1.25)
tbsh.width = Inches(12.5); tbsh.height = Inches(1.55)
tbl = tbsh.table
colw = [1.5, 1.9, 1.4, 1.8, 1.2, 1.9, 2.8]
for i, w in enumerate(colw):
    tbl.columns[i].width = Inches(w)
hdr_cells = ["Subject", "Loss combo", "RDM L_test", "RDM ΔL vs(0,0)",
             "folds<(0,0)", "γ ΔL vs(0,0)", "Read"]
for j, t in enumerate(hdr_cells):
    set_cell(tbl.cell(0, j), t, size=11, bold=True, col=WHITE)
row1 = ["S08 (6,−42)", "γOY+RDM_V2", "0.594", "−0.406 ✓", "7/7", "−13.8 (5/7)", "neural + behav"]
row2 = ["S09 (2,+24)", "γall+RDM_V1", "0.528", "−0.472 ✓", "7/7", "−0.55 (4/7)", "neural-only"]
for j, t in enumerate(row1):
    set_cell(tbl.cell(1, j), t, size=11, bold=False, col=DARK)
for j, t in enumerate(row2):
    set_cell(tbl.cell(2, j), t, size=11, bold=False, col=DARK)
# row3: merged caveat
tbl.cell(3, 0).merge(tbl.cell(3, 6))
set_cell(tbl.cell(3, 0),
         "(0,0) = no-correction floor (RDM≡1.0).  ΔL<0 = beats it on EVERY held-out HC.  "
         "grid pct 92–95% (non-trivial, not just the floor).  lower L_test = better.",
         size=10, bold=False, col=GRAY)

# left card (chain): box[7] keep; header text[8]; body text[9]
shp(ns, 7).top = Inches(3.05); shp(ns, 7).height = Inches(3.75)
shp(ns, 8).top = Inches(3.13)
rebuild_textbox(shp(ns, 8),
                [("왜 test-loss인가 — stability ≠ goodness", 15, True, DARK)])
b9 = shp(ns, 9); b9.top = Inches(3.55); b9.height = Inches(3.15)
rebuild_textbox(b9, [
    ("Stability (s17): HC 부분집합 refit 재현성 = estimator 분산.", 13, False, DARK),
    ("  → misspecified여도 같은 값 수렴 가능; 값의 우수성·overfitting 미검출.", 12, False, GRAY),
    ("Held-out test-loss (s18): 6 HC fit → 본 적 없는 HC 예측.", 13, False, DARK),
    ("  → overfitting + 임의성 배제. stability 넘어서는 평가 기준 (e).", 12, False, GREEN),
    ("두 term 모두 무보정(δθ=0) 대비 ΔL로 채점 (uniform):", 13, False, DARK),
    ("  γ = held-out HC를 baseline, target=CVD JND (reference-robust).", 12, False, GRAY),
    ("  RDM = held-out HC 기하가 target (genuine prediction).", 12, False, GRAY),
])

# right card (reconciliation): box[10] -> move right; text[11] header; text[12] body
shp(ns, 10).left = Inches(6.85); shp(ns, 10).top = Inches(3.05)
shp(ns, 10).width = Inches(6.10); shp(ns, 10).height = Inches(3.75)
h11 = shp(ns, 11); h11.left = Inches(7.00); h11.top = Inches(3.13); h11.width = Inches(5.85)
rebuild_textbox(h11, [("해석 — 하나의 broad shallow basin", 15, True, DARK)])
b12 = shp(ns, 12); b12.left = Inches(7.00); b12.top = Inches(3.55)
b12.width = Inches(5.85); b12.height = Inches(3.15)
rebuild_textbox(b12, [
    ("neural 안정값: 두 피험자 모두 무보정보다 우수 (7/7 fold).", 13, False, GREEN),
    ("behav: S08 강함 (triangulation) / S09 ≈ null.", 13, False, DARK),
    ("Test 2a ~20° non-identifiability와 모순 없음 — 같은 basin:", 13, False, DARK),
    ("  중심 일관 (stable)", 12, False, GRAY),
    ("  공유 + 무보정 우위 (good)", 12, False, GRAY),
    ("  폭 ~20° (절대값 미고정)", 12, False, AMBER),
    ("→ 값은 'good region', 점-정밀 아님 (§0 descriptive, specificity 아님).", 13, True, DARK),
])

# ---------------------------------------------------------------- slide 7 (RQ4)
s7 = find_slide_by_title("RQ4+RQ5")
# S09 critical-finding box (shape id 11): add OOS-good + honest caveat
s7_11 = shp(s7, 11)
append_cloned(s7_11, [
    (2, "OOS: 무보정 대비 7/7 fold 우위 (s18) ✓"),     # clone green bold p2
    (3, "⚠ but behav≈0 · 값 ~20° basin 미고정 (not pinned)"),  # clone normal p3
])
# bottom justification box (id 16): add triangulation balance line
s7_16 = shp(s7, 16)
append_cloned(s7_16, [
    (2, "S08: behav+neural 같은 방향 (triangulation); S09: neural-only (behav silent)"),
])

# ---------------------------------------------------------------- slide 8 (RQ3)
s8 = find_slide_by_title("RQ3 — Identifiability")
# find the largest text box to append the basin line; use the box containing 'Descriptive'
target8 = None
for sh in s8.shapes:
    if sh.has_text_frame and "Descriptive embedding" in sh.text_frame.text:
        target8 = sh; break
if target8 is None:  # fallback: last big text box
    target8 = max((sh for sh in s8.shapes if sh.has_text_frame),
                  key=lambda x: (x.width or 0) * (x.height or 0))
np_ = len(target8.text_frame.paragraphs)
append_cloned(target8, [
    (np_ - 1, "Held-out test-loss + stability + ~20° floor = one broad shallow "
              "basin (centered · shared · wide) — s18"),
])

# ---------------------------------------------------------------- slide 12 (limits)
s12 = find_slide_by_title("한계 및 다음 단계")
for sh in s12.shapes:
    if not sh.has_text_frame:
        continue
    for para in sh.text_frame.paragraphs:
        txt = "".join(r.text for r in para.runs)
        if "All OOS axes = HC pool only" in txt:
            set_run_text(para._p, "HC-pool OOS: held-out test-loss ✓ (s18, 7/7)")
        elif txt.strip() == "→ Generalization = Phase 3":
            set_run_text(para._p, "→ CVD-level OOS = Phase 3 (N=2)")

# ---------------------------------------------------------------- renumber pages
for i, s in enumerate(prs.slides, start=1):
    for sh in s.shapes:
        if (sh.has_text_frame and sh.left is not None
                and Inches(12.2) <= sh.left <= Inches(13.0)
                and sh.top is not None and sh.top <= Inches(0.5)):
            t = sh.text_frame.text.strip()
            if t.isdigit():
                for para in sh.text_frame.paragraphs:
                    if any(r.text.strip().isdigit() for r in para.runs):
                        set_run_text(para._p, str(i))

prs.save(SRC)
print("saved", SRC, "| n_slides =", len(prs.slides))
