#!/usr/bin/env python3
"""
patch_fig3_box2.py

Replaces box 2 of the pipeline schematic (`fig:pipeline`, file stem
`fig3_workflow`) inside the PowerPoint source, then leaves a .pptx ready for
export.

`fig3_workflow` is the one manuscript figure that is a hand-assembled
PowerPoint composite rather than a single script.  Its slide carries a
full-slide background PICTURE that is a rendering of an OLDER version of the
whole figure, with live shapes overlaid on top.  Deleting the box-2 shapes
therefore is not enough: the old box-2 content would show through from that
background.  The new asset is opaque and is placed over the box-2 interior,
which both hides the background and supplies the new content.

What is removed (2026-09-03, author instruction; see
generate_box2_loss_gates.py for the rationale):
    id 17  bullet list of the three loss atoms
    id 20  "not selected" tag
    id 25  "3-gate selection (...)" text
    id 26  the grey box behind it
    id 27  "deutan: L_gamma(OY) + L_RDM(V2)"
    id 28  "protan: L_gamma(all) + L_RDM(V1)"
    id 30  connector to id 27
    id 31  connector to id 28

Geometry: the box-2 frame is part of the background picture.  Its interior was
measured off that picture (137.22 px/in horizontally, 155.60 px/in vertically,
picture anchored at 0.03 in, 1.58 in) and is 3.18-5.69 in x 2.14-6.40 in.

The embedded hue-wheel SVGs (`box1_forward_wheel`, `box4_preimage_wheel`)
declare matplotlib's default font stack, which lists "DejaVu Sans" ahead of
Arial.  PowerPoint happened to resolve that to Arial, but any renderer with
DejaVu installed picks DejaVu and the wheel labels stop being Arial, which the
journal requires (Figures/scripts/FONT_POLICY.md).  The font stack is therefore
rewritten in place; this is what "arial_fontpatched" already did for the slide
text but not for the SVG interiors.

Usage:
    python patch_fig3_box2.py            # writes the patched .pptx
    python patch_fig3_box2.py --export   # ... and exports a PDF via LibreOffice
"""
import re
import shutil
import subprocess
import sys
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

ASSETS = Path(__file__).resolve().parent.parent / "fig3_assets"
SRC = ASSETS / "Presentation1_arial_fontpatched.pptx"
DST = ASSETS / "Presentation1_box2_2026-09-03.pptx"
NEW_BOX2 = ASSETS / "box2_loss_gates.png"

DROP_IDS = {17, 20, 25, 26, 27, 28, 30, 31}
BOX2 = dict(left=3.18, top=2.14, width=2.51, height=4.26)   # inches

SOFFICE = "/Applications/LibreOffice.app/Contents/MacOS/soffice"


ARIAL_STACK = "Arial,Helvetica,sans-serif"


def _pin_svg_font(pptx_path):
    """Force the embedded SVGs onto Arial (see the module docstring)."""
    with zipfile.ZipFile(pptx_path) as z:
        items = [(i, z.read(i.filename)) for i in z.infolist()]
    patched = []
    tmp = pptx_path.with_suffix(".tmp.pptx")
    with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as z:
        for info, data in items:
            if info.filename.lower().endswith(".svg"):
                text = data.decode("utf-8")
                new = re.sub(r'font-family="[^"]*"',
                             f'font-family="{ARIAL_STACK}"', text)
                if new != text:
                    patched.append(info.filename)
                data = new.encode("utf-8")
            z.writestr(info, data)
    tmp.replace(pptx_path)
    print("font-pinned SVGs:", patched)


def main():
    if not NEW_BOX2.exists():
        sys.exit(f"missing asset: {NEW_BOX2} (run generate_box2_loss_gates.py)")
    shutil.copy(SRC, DST)

    prs = Presentation(DST)
    slide = prs.slides[0]

    removed = []
    for shape in list(slide.shapes):
        if shape.shape_id in DROP_IDS:
            shape._element.getparent().remove(shape._element)
            removed.append(shape.shape_id)
    missing = DROP_IDS - set(removed)
    if missing:
        sys.exit(f"expected shapes not found, aborting: {sorted(missing)}")
    print("removed shape ids:", sorted(removed))

    slide.shapes.add_picture(
        str(NEW_BOX2),
        Inches(BOX2["left"]), Inches(BOX2["top"]),
        Inches(BOX2["width"]), Inches(BOX2["height"]),
    )
    prs.save(DST)
    _pin_svg_font(DST)
    print("saved:", DST)

    if "--export" in sys.argv:
        subprocess.run([SOFFICE, "--headless", "--convert-to", "pdf",
                        "--outdir", str(ASSETS), str(DST)], check=True)
        print("exported:", ASSETS / (DST.stem + ".pdf"))


if __name__ == "__main__":
    main()
