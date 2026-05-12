"""
Build a PowerPoint (16:9) summarising sub-08 filter candidates with editable text.

Prerequisites:
    pip install python-pptx pillow numpy

Usage:
    python scripts/pptx_candidates.py
Output:
    candidates/filter_candidates_sub08.pptx   (editable in PowerPoint / Keynote)

Slide layout (per filter):
    ┌─────────────────────────────┬─────────────────────────────┐
    │  cropped c1–c8 image        │  Behavioral table           │
    │  (4 visual columns)         │  Color | Original | Filtered │
    │                             │  c1 …                       │
    │                             │  c8 …                       │
    └─────────────────────────────┴─────────────────────────────┘
"""

import os
import sys
import textwrap

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
    import numpy as np
    from PIL import Image
except ImportError:
    sys.exit("Missing dependency — run:  pip install python-pptx pillow numpy")

# ── Behavioral data ──────────────────────────────────────────────────────────

COLOR_LABELS = {
    "c1": "C1 magenta-red (0°)",
    "c2": "C2 orange (45°)",
    "c3": "C3 yellow (90°)",
    "c4": "C4 yellow-green (135°)",
    "c5": "C5 teal/green (180°)",
    "c6": "C6 cyan-blue (225°)",
    "c7": "C7 blue (270°)",
    "c8": "C8 purple (315°)",
}

ORIGINAL_REPORT = {
    "c1": "빨강에 분홍 섞인",
    "c2": "연한 빨강",
    "c3": "연두",
    "c4": "노랑",
    "c5": "웜톤 아이보리",
    "c6": "하늘",
    "c7": "더 진한 하늘",
    "c8": "진파랑",
}

FILTERS = [
    {
        "key":   "canonical",
        "title": "Canonical  —  β_s = 38°,  β_c = −14°",
        "note":  "hV4 LOCO ρ argmax · behavioral PASS (2026-04-17)",
        "data": {
            "c1": ("분홍기↓ 순수 빨강 가까워짐", "pass"),
            "c2": ("약간 초록 느낌, 원본보다 옅음",  "note"),
            "c3": ("C4와 같은 노란색  [collapse]",   "fail"),
            "c4": ("C3와 같은 노란색  [collapse]",   "fail"),
            "c5": ("C6와 같은 하늘색  [swap]",       "fail"),
            "c6": ("C5와 같은 하늘, 원본 C7  [swap]","fail"),
            "c7": ("원본보다 더 진한 파란색",          "note"),
            "c8": ("원본과 같음 (무변화)",             "note"),
        },
    },
    {
        "key":   "primary",
        "title": "V4-only  —  β_s = 38°,  β_c = +7°",
        "note":  "z_combined CI disjoint from HC  (β_c sign reversed vs Canonical)",
        "data": {
            "c1": ("(언급 없음)",          "note"),
            "c2": ("더 진한 빨강",          "note"),
            "c3": ("연한 주황",             "note"),
            "c4": ("연한 연두",             "note"),
            "c5": ("완전 옅은 연두",        "note"),
            "c6": ("더 진한 하늘색",        "note"),
            "c7": ("완전 짙은 파랑",        "note"),
            "c8": ("핑크빛 섞인 보라",      "note"),
        },
    },
    {
        "key":   "cycle14",
        "title": "Cycle 14  —  β_s = 58°,  β_c = −36°",
        "note":  "2·mwJ(V4) + 1·rdm(V1) + 0.2·Tikh · boot_frac = 1.000 (strongest neural)",
        "data": {
            "c1": ("보라에서 핑크  [역방향]",         "fail"),
            "c2": ("더 옅은 주황",                   "note"),
            "c3": ("연한 주황에서 노랑",              "note"),
            "c4": ("연둣빛 노랑 → 웜톤 아이보리",    "note"),
            "c5": ("웜톤 아이보리에서 하늘",          "note"),
            "c6": ("약간 진한 하늘",                  "note"),
            "c7": ("파랑에서 하늘",                   "note"),
            "c8": ("진한 파랑에서 파랑",              "note"),
        },
    },
]

# ── Colour helpers ───────────────────────────────────────────────────────────

PASS_RGB = RGBColor(0x00, 0x8B, 0x00)
FAIL_RGB = RGBColor(0xB0, 0x00, 0x00)
NOTE_RGB = RGBColor(0x30, 0x30, 0x30)
HDR_RGB  = RGBColor(0x1E, 0x1E, 0x5A)

STATUS_COLOR = {"pass": PASS_RGB, "fail": FAIL_RGB, "note": NOTE_RGB}

# ── Image cropping ───────────────────────────────────────────────────────────

def detect_tier1_bounds(img_array, probe_x1=260, probe_x2=380):
    strip      = img_array[:, probe_x1:probe_x2, :3].astype(float)
    brightness = np.mean(strip, axis=(1, 2))
    std        = np.std(strip, axis=(1, 2))
    segs, state, seg_start = [], "white", 150
    for y in range(150, min(1200, img_array.shape[0])):
        cur = "white" if brightness[y] > 245 and std[y] < 15 else "color"
        if cur != state:
            segs.append((state, seg_start, y - 1))
            state, seg_start = cur, y
    segs.append((state, seg_start, min(1199, img_array.shape[0] - 1)))
    color_segs = [(y0, y1) for kind, y0, y1 in segs if kind == "color" and (y1 - y0) > 30]
    header_seg = color_segs[0] if color_segs and (color_segs[0][1] - color_segs[0][0]) < 35 else None
    main_segs  = [s for s in color_segs if s is not header_seg][:8]
    if not main_segs:
        return 150, min(900, img_array.shape[0] - 1)
    y_top    = (header_seg[0] if header_seg else main_segs[0][0]) - 10
    y_bottom = main_segs[-1][1] + 10
    return max(0, y_top), min(img_array.shape[0] - 1, y_bottom)


def cropped_png_path(src_path: str) -> str:
    """Return a temp cropped PNG path (next to the source)."""
    d, name = os.path.split(src_path)
    base, _ = os.path.splitext(name)
    out = os.path.join(d, f"_tmp_cropped_{base}.png")
    img = Image.open(src_path).convert("RGB")
    arr = np.array(img)
    y0, y1 = detect_tier1_bounds(arr)
    img.crop((0, y0, img.width, y1)).save(out, dpi=(150, 150))
    return out


# ── Slide builder ────────────────────────────────────────────────────────────

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

IMG_LEFT   = Inches(0.2)
IMG_TOP    = Inches(1.0)
IMG_W      = Inches(8.2)          # image panel width
TBL_LEFT   = Inches(8.6)
TBL_TOP    = Inches(1.0)
TBL_W      = Inches(4.5)
TBL_H      = Inches(6.2)


def set_cell_text(cell, text: str, bold=False, font_size=Pt(11),
                  color: RGBColor = NOTE_RGB, align=PP_ALIGN.LEFT):
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color


def set_cell_bg(cell, r, g, b):
    from pptx.oxml.ns import qn
    from lxml import etree
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", f"{r:02X}{g:02X}{b:02X}")


def add_filter_slide(prs: Presentation, filt: dict, img_path: str):
    slide_layout = prs.slide_layouts[6]   # blank layout
    slide = prs.slides.add_slide(slide_layout)

    # ── Title text box ────────────────────────────────────────────────────────
    txb = slide.shapes.add_textbox(Inches(0.2), Inches(0.05), Inches(13.0), Inches(0.85))
    tf  = txb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = filt["title"]
    run.font.size  = Pt(24)
    run.font.bold  = True
    run.font.color.rgb = HDR_RGB

    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = filt["note"]
    r2.font.size  = Pt(13)
    r2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # ── Cropped image ────────────────────────────────────────────────────────
    crop_path = cropped_png_path(img_path)
    pic = slide.shapes.add_picture(crop_path, IMG_LEFT, IMG_TOP, width=IMG_W)
    # Constrain height to slide area
    max_h = SLIDE_H - IMG_TOP - Inches(0.1)
    if pic.height > max_h:
        ratio    = max_h / pic.height
        pic.height = int(max_h)
        pic.width  = int(pic.width * ratio)

    # ── Behavioral table ─────────────────────────────────────────────────────
    # Columns: Color label | Original (CVD sees) | Filtered (CVD sees)
    rows = 9   # 1 header + 8 color rows
    cols = 3
    tbl  = slide.shapes.add_table(rows, cols, TBL_LEFT, TBL_TOP, TBL_W, TBL_H).table

    # Column widths
    tbl.columns[0].width = Inches(1.55)
    tbl.columns[1].width = Inches(1.45)
    tbl.columns[2].width = Inches(1.5)

    # Header row
    headers = ["Color", "Original\n(sub-08)", "After Filter\n(sub-08)"]
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        set_cell_text(cell, hdr, bold=True, font_size=Pt(11),
                      color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
        set_cell_bg(cell, 0x1E, 0x1E, 0x5A)

    # Data rows
    row_bg_even = (0xF0, 0xF4, 0xFF)
    row_bg_odd  = (0xFF, 0xFF, 0xFF)
    for ri, ck in enumerate(["c1","c2","c3","c4","c5","c6","c7","c8"]):
        row_idx   = ri + 1
        bg        = row_bg_even if ri % 2 == 0 else row_bg_odd
        lbl       = COLOR_LABELS[ck]
        orig_txt  = ORIGINAL_REPORT[ck]
        filt_txt, status = filt["data"][ck]

        for ci in range(3):
            set_cell_bg(tbl.cell(row_idx, ci), *bg)

        set_cell_text(tbl.cell(row_idx, 0), lbl,      font_size=Pt(10), color=NOTE_RGB)
        set_cell_text(tbl.cell(row_idx, 1), orig_txt, font_size=Pt(10), color=NOTE_RGB)
        set_cell_text(tbl.cell(row_idx, 2), filt_txt,
                      font_size=Pt(10), color=STATUS_COLOR[status],
                      bold=(status != "note"))

    return slide


def add_title_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf  = txb.text_frame
    p   = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Sub-08 Deutan — Filter Candidates\nBehavioral Validation Comparison"
    run.font.size  = Pt(36)
    run.font.bold  = True
    run.font.color.rgb = HDR_RGB
    p.alignment = PP_ALIGN.CENTER

    p2  = tf.add_paragraph()
    r2  = p2.add_run()
    r2.text = "2-Component model  ·  CIELab opponent space  ·  hV4 LOCO primary criterion"
    r2.font.size  = Pt(18)
    r2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    p2.alignment = PP_ALIGN.CENTER


# ── Main ─────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    base = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    cand = os.path.join(base, "candidates")

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    add_title_slide(prs)

    for filt in FILTERS:
        src = os.path.join(cand, f"{filt['key']}.png")
        if not os.path.exists(src):
            print(f"  Skipping {src} (not found)")
            continue
        print(f"  Adding slide: {filt['title']}")
        add_filter_slide(prs, filt, src)

    out = os.path.join(cand, "filter_candidates_sub08.pptx")
    prs.save(out)
    print(f"\nSaved: {out}")

    # Clean up temp files
    for filt in FILTERS:
        tmp = os.path.join(cand, f"_tmp_cropped_{filt['key']}.png")
        if os.path.exists(tmp):
            os.remove(tmp)
