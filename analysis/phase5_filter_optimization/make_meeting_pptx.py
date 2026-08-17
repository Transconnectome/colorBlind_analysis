#!/usr/bin/env python3
"""
Meeting presentation — Phase 2 CVD color filter.
All text trimmed to fit boxes. Min 14pt body (13pt gray sub-notes).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT     = "meeting_2026-06-01.pptx"
FIG_DIR = "results/figures"
V4COL   = "results/visualizations/pipeline2_primary_4col"

NAVY    = RGBColor(0x0D, 0x2B, 0x4F)
BLUE_H  = RGBColor(0x1A, 0x56, 0x8C)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BODY    = RGBColor(0x1A, 0x1A, 0x1A)
LBLUE   = RGBColor(0xE8, 0xF2, 0xFB)
S08C    = RGBColor(0xBE, 0x35, 0x25)
S09C    = RGBColor(0x0E, 0x6B, 0x4F)
AMBER   = RGBColor(0xD4, 0x7E, 0x00)
GREEN   = RGBColor(0x1E, 0x8B, 0x4C)
REDF    = RGBColor(0xC0, 0x29, 0x1E)
GRAY    = RGBColor(0x7F, 0x8C, 0x8D)
DIVIDER = RGBColor(0xD5, 0xDB, 0xE0)
PURPLE  = RGBColor(0x7B, 0x26, 0x8C)
LIGHT_R = RGBColor(0xFD, 0xF3, 0xF2)
LIGHT_G = RGBColor(0xF0, 0xF9, 0xF4)

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def box(slide, l, t, w, h, fill=None, border=None, bw=1):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border; s.line.width = Pt(bw)
    else:
        s.line.fill.background()
    return s

def tb(slide, l, t, w, h, wrap=True):
    b = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    b.text_frame.word_wrap = wrap
    return b.text_frame

def p(tf, text, sz=16, bold=False, col=BODY, align=PP_ALIGN.LEFT, sb=0, italic=False):
    pg = tf.add_paragraph()
    pg.alignment = align; pg.space_before = Pt(sb)
    r = pg.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = col
    return pg

def bl(tf, text, sz=14, col=BODY, sb=1, bold=False):
    pg = tf.add_paragraph()
    pg.alignment = PP_ALIGN.LEFT; pg.space_before = Pt(sb)
    r = pg.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = col
    return pg

def hdr(slide, title, n, sub=None):
    box(slide, 0, 0, 13.33, 1.05, fill=NAVY)
    t = tb(slide, 0.35, 0.08, 11.7, 0.82)
    p(t, title, sz=22, bold=True, col=WHITE)
    nt = tb(slide, 12.5, 0.08, 0.7, 0.35, wrap=False)
    p(nt, str(n), sz=12, col=GRAY, align=PP_ALIGN.RIGHT)
    if sub:
        st = tb(slide, 0.35, 0.82, 11.6, 0.27)
        p(st, sub, sz=12, col=RGBColor(0xA8, 0xCC, 0xEC), italic=True)

def add_table(slide, hdrs, rows, l, t, w, h, hcol=BLUE_H):
    nr, nc = len(rows)+1, len(hdrs)
    tbl = slide.shapes.add_table(nr, nc, Inches(l), Inches(t), Inches(w), Inches(h)).table
    for j, hd in enumerate(hdrs):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = hcol
        pg = c.text_frame.paragraphs[0]; pg.alignment = PP_ALIGN.CENTER
        r = pg.add_run(); r.text = hd
        r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = WHITE
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = tbl.cell(i+1, j)
            bg = RGBColor(0xF5,0xF8,0xFF) if i%2==0 else WHITE
            c.fill.solid(); c.fill.fore_color.rgb = bg
            pg = c.text_frame.paragraphs[0]; pg.alignment = PP_ALIGN.CENTER
            r = pg.add_run(); r.text = str(val)
            r.font.size = Pt(12); r.font.color.rgb = BODY
    return tbl

def img(slide, path, l, t, w, h=None):
    if not os.path.exists(path): return
    from PIL import Image as PI
    iw, ih = PI.open(path).size
    height = Inches(h) if h else Inches(w * ih / iw)
    slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), height)

# ═══════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── SLIDE 1: TITLE ─────────────────────────────────────────────
sl = blank(prs)
box(sl, 0, 0, 13.33, 7.5, fill=NAVY)
box(sl, 0, 5.5, 13.33, 2.0, fill=BLUE_H)
box(sl, 0, 5.47, 13.33, 0.08, fill=RGBColor(0x4A,0x90,0xD9))

tf = tb(sl, 1.0, 1.1, 11.3, 1.8)
p(tf, "CVD 색 인지 보정 필터", sz=30, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
p(tf, "Personalized Filter Design via Neural Inverse Fitting",
  sz=20, col=RGBColor(0xA8,0xCC,0xEC), align=PP_ALIGN.CENTER, sb=6)

tf2 = tb(sl, 1.0, 3.1, 11.3, 1.5)
p(tf2, "Phase 2: Model Selection · Fitting · Pre-image Filter",
  sz=17, col=RGBColor(0xCC,0xDD,0xEE), align=PP_ALIGN.CENTER)
p(tf2, "Sub-08 (Deutan)  ·  Sub-09 (Protan)",
  sz=16, col=RGBColor(0xCC,0xDD,0xEE), align=PP_ALIGN.CENTER, sb=4)

tf3 = tb(sl, 1.0, 5.65, 11.3, 0.6)
p(tf3, "2026-06-01  ·  미팅 발표자료",
  sz=14, col=RGBColor(0x88,0xAA,0xCC), align=PP_ALIGN.CENTER)


# ── SLIDE 2: PHASE 2 PIPELINE (5-STEP) ────────────────────────
sl = blank(prs)
hdr(sl, "Phase 2 파이프라인 (5-Step Selection Axis)", 2,
    sub="Models · Loss atoms · Fitting · Stability · Final decision + Limits")

SCOLS = [BLUE_H, RGBColor(0x2E,0x86,0xAB), RGBColor(0x0E,0x6B,0x4F),
         PURPLE, REDF]
BW = 2.44; GAP = 0.07; LX0 = 0.38

for i in range(5):
    lx  = LX0 + i*(BW+GAP)
    col = SCOLS[i]
    box(sl, lx, 1.15, BW, 0.52, fill=col)
    th = tb(sl, lx+0.08, 1.18, BW-0.16, 0.46)
    p(th, f"Step {i+1}", sz=13, bold=True, col=WHITE)
    box(sl, lx, 1.67, BW, 5.6, fill=WHITE, border=col, bw=1)
    if i < 4:
        ta = tb(sl, lx+BW+0.005, 4.2, 0.075, 0.4)
        p(ta, "›", sz=18, bold=True, col=col, align=PP_ALIGN.CENTER)

# Step content — each line ≤ ~24 chars at 13pt in 2.44" col
def sc(slide, i, lines):
    """Add step content: list of (text, size, color, bold, space_before)."""
    lx = LX0 + i*(BW+GAP)
    tf2 = tb(slide, lx+0.1, 1.72, BW-0.2, 5.5)
    for text, sz, col, bold, sb in lines:
        pg = tf2.add_paragraph()
        pg.space_before = Pt(sb)
        r = pg.add_run(); r.text = text
        r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = col

sc(sl, 0, [
    ("Models",         14, BLUE_H, True,  0),
    ("R+C",            13, BODY,   False, 3),
    ("δθ=(2−g)·δθ_M", 13, GRAY,   False, 0),
    ("1 DOF  (g)",     13, GRAY,   False, 0),
    ("2-Component",    13, BODY,   False, 5),
    ("δθ=β_s·cos(…)",  13, GRAY,   False, 0),
    ("     +β_c·cos(…)",13,GRAY,   False, 0),
    ("2 DOF  (β_s,β_c)",13,GRAY,  False, 0),
    ("Loss Atoms",     14, BLUE_H, True,  8),
    ("γ — JND z²",    13, BODY,   False, 3),
    ("  per-pair + all",13,GRAY,   False, 0),
    ("L_RDM — PCA 45°",13,BODY,   False, 3),
    ("  28-d cosine",  13, GRAY,   False, 0),
    ("L_LOCO — hV4",   13, BODY,   False, 3),
    ("  voxel predict.",13,GRAY,   False, 0),
    ("  (only fwd-mdl)",13,BLUE_H, False, 0),
])

sc(sl, 1, [
    ("Cell Enum.",     14, RGBColor(0x2E,0x86,0xAB), True,  0),
    ("S08: 71×4 cells",13, BODY,   False, 3),
    ("S09: 11×4 cells",13, BODY,   False, 0),
    ("",               12, BODY,   False, 2),
    ("Precond. gate",  14, RGBColor(0x2E,0x86,0xAB), True,  5),
    ("d ≥ +0.5",       13, BODY,   False, 3),
    ("(CVD > HC dir.)",13, GRAY,   False, 0),
    ("",               12, BODY,   False, 2),
    ("Pass rate",      14, RGBColor(0x2E,0x86,0xAB), True,  5),
    ("S08: 31/284 (11%)",13,BODY,  False, 3),
    ("S09:  2/44   (5%)",13,BODY,  False, 0),
])

sc(sl, 2, [
    ("HC Resample",    14, RGBColor(0x0E,0x6B,0x4F), True,  0),
    ("5-train/2-test", 13, BODY,   False, 3),
    ("N=300, seed=42", 13, GRAY,   False, 0),
    ("",               12, BODY,   False, 2),
    ("Composite fit",  14, RGBColor(0x0E,0x6B,0x4F), True,  5),
    ("z_sum=Σzscore/√n",13,BODY,   False, 3),
    ("argmin → (β_s,β_c)",13,GRAY, False, 0),
    ("σ-bin: 45° RDM", 13, GRAY,   False, 0),
    ("",               12, BODY,   False, 2),
    ("Strict HC LOO",  14, RGBColor(0x0E,0x6B,0x4F), True,  5),
    ("7-fold (s17)",   13, BODY,   False, 3),
    ("deterministic",  13, GRAY,   False, 0),
])

sc(sl, 3, [
    ("Stability",      14, PURPLE, True,  0),
    ("Boundary < 50%", 13, BODY,   False, 3),
    ("param IQR ↓",    13, BODY,   False, 0),
    ("mode share ↑",   13, BODY,   False, 0),
    ("LOO range ↓",    13, BODY,   False, 0),
    ("",               12, BODY,   False, 3),
    ("Top candidates", 14, PURPLE, True,  5),
    ("S08 βc-dom",     13, S08C,   True,  3),
    ("  (6, −42)",     13, S08C,   False, 0),
    ("  IQR(8,2) ~70%",13, BODY,   False, 0),
    ("S09 βc-rot",     13, S09C,   True,  3),
    ("  (2, +24)",     13, S09C,   False, 0),
    ("  IQR(0,0) 87.7%",13,BODY,  False, 0),
])

sc(sl, 4, [
    ("Signal check",   14, GREEN,  True,  0),
    ("Test loss depth",13, BODY,   False, 3),
    ("REAL vs SYNTH:", 13, BODY,   False, 0),
    ("S08: 5.5× deeper",13,GREEN,  True,  0),
    ("S09: 3.9× deeper",13,GREEN,  True,  0),
    ("→ signal used ✓",13, GREEN,  False, 0),
    ("",               12, BODY,   False, 2),
    ("Specificity",    14, REDF,   True,  3),
    ("T1 f10°<0.30 FAIL",13,REDF,  False, 3),
    ("T2a 0/140  FAIL",13, REDF,   False, 0),
    ("T2b/2c  0/3 FAIL",13,REDF,   False, 0),
    ("Pre-image",      14, GREEN,  True,  5),
    ("8/8 exact ✓",    13, GREEN,  True,  3),
    ("★ Descriptive",  13, AMBER,  True,  5),
    ("  embed. only",  13, AMBER,  False, 0),
])


# ── SLIDE 3: TWO MODELS ────────────────────────────────────────
sl = blank(prs)
hdr(sl, "두 후보 모델 (M1)", 3,
    sub="R+C (1 DOF) vs 2-Component (2 DOF)  —  neutral in Methods; adequacy judged in Results")

box(sl, 0.4, 1.2, 5.9, 5.6, fill=RGBColor(0xFE,0xF9,0xF0), border=AMBER, bw=2)
t = tb(sl, 0.6, 1.3, 5.5, 0.4)
p(t, "Candidate 1 — R+C", sz=18, bold=True, col=AMBER)
t2 = tb(sl, 0.6, 1.75, 5.5, 0.75)
p(t2, "δθ_RC = (2 − g) · δθ_Machado(c; Δλ)", sz=15, bold=True, col=NAVY, align=PP_ALIGN.CENTER)
p(t2, "g=1: none  |  g=2: perfect  |  g>2: over", sz=13, col=GRAY, align=PP_ALIGN.CENTER, sb=2)
ln = sl.shapes.add_connector(1, Inches(0.55), Inches(2.58), Inches(6.25), Inches(2.58))
ln.line.color.rgb = DIVIDER; ln.line.width = Pt(1)
t3 = tb(sl, 0.6, 2.65, 5.5, 3.9)
p(t3,  "Mechanism",       sz=14, bold=True, col=AMBER)
bl(t3, "Retinal cone shift (Δλ, fixed)", sz=14)
bl(t3, "Cortical gain g  [1 free param]", sz=14)
p(t3,  "Parameter space", sz=14, bold=True, col=AMBER, sb=6)
bl(t3, "g ∈ [1.0, 3.0], step 0.1  (DOF=1)", sz=14)
bl(t3, "Δλ fixed: DPS_lit / Boehm / JND", sz=14)
p(t3,  "Literature",      sz=14, bold=True, col=AMBER, sb=6)
bl(t3, "Machado 2009  (cone shift)", sz=13, col=GRAY)
bl(t3, "Boehm 2014, Tregillus 2021", sz=13, col=GRAY)

box(sl, 7.0, 1.2, 5.9, 5.6, fill=RGBColor(0xF0,0xF6,0xFF), border=BLUE_H, bw=2)
t = tb(sl, 7.2, 1.3, 5.5, 0.4)
p(t, "Candidate 2 — 2-Component  ★", sz=18, bold=True, col=BLUE_H)
t2 = tb(sl, 7.2, 1.75, 5.5, 0.75)
p(t2, "δθ = β_s·cos(θ−90°) + β_c·cos(θ−θ_conf)", sz=15, bold=True, col=NAVY, align=PP_ALIGN.CENTER)
p(t2, "θ_conf: protan=16°,  deutan=150°", sz=13, col=GRAY, align=PP_ALIGN.CENTER, sb=2)
ln2 = sl.shapes.add_connector(1, Inches(7.15), Inches(2.58), Inches(12.85), Inches(2.58))
ln2.line.color.rgb = DIVIDER; ln2.line.width = Pt(1)
t3 = tb(sl, 7.2, 2.65, 5.5, 3.9)
p(t3,  "Mechanism",       sz=14, bold=True, col=BLUE_H)
bl(t3, "β_s: S-cone cardinal axis rotation", sz=14)
bl(t3, "β_c: confusion-axis rotation", sz=14)
bl(t3, "β_s ≥ 0  (Emery 2021/23 direction)", sz=14)
p(t3,  "Parameter space", sz=14, bold=True, col=BLUE_H, sb=6)
bl(t3, "β_s ∈ [0°,50°], β_c ∈ [−50°,50°]", sz=14)
bl(t3, "step 2°,  DOF=2", sz=14)
p(t3,  "Novel formulation", sz=14, bold=True, col=BLUE_H, sb=6)
bl(t3, "No prior art — cardinal convention only", sz=13, col=GRAY)
bl(t3, "NOT derived from existing model", sz=13, col=GRAY)

tf_vs = tb(sl, 6.1, 3.6, 0.8, 0.8)
p(tf_vs, "vs", sz=26, bold=True, col=NAVY, align=PP_ALIGN.CENTER)


# ── SLIDE 4: INVERSE FITTING + CRITERIA ───────────────────────
sl = blank(prs)
hdr(sl, "역문제 피팅 + 모델 선택 기준 (M2 + M2.5)", 4,
    sub="Inverse problem: observed neural & behavioral signals → (β_s, β_c)")

box(sl, 0.4, 1.2, 5.85, 5.5, fill=LBLUE, border=BLUE_H, bw=1)
t = tb(sl, 0.55, 1.3, 5.55, 0.4)
p(t, "Three Loss Families", sz=17, bold=True, col=BLUE_H)
t2 = tb(sl, 0.55, 1.75, 5.55, 4.8)
p(t2,  "① Behavioral JND  (γ)",    sz=15, bold=True, col=BODY)
bl(t2, "Per-pair z²: γ_OY/YG/YP/GB", sz=14)
bl(t2, "Aggregate γ_all (8-pair sum)", sz=14)
bl(t2, "No forward model needed",    sz=13, col=GRAY)
p(t2,  "② Repr. RDM  (L_RDM)",     sz=15, bold=True, col=BODY, sb=8)
bl(t2, "PCA top-6 → 8×8 RDM",       sz=14)
bl(t2, "45° σ-bin, cosine to HC mean",sz=14)
bl(t2, "No forward model needed",    sz=13, col=GRAY)
p(t2,  "③ LOCO voxel  (L_LOCO)",   sz=15, bold=True, col=BODY, sb=8)
bl(t2, "hV4 voxel prediction",       sz=14)
bl(t2, "ridge-GCV encoder",          sz=14)
bl(t2, "★ Only atom needing fwd model",sz=14, bold=True, col=BLUE_H)

box(sl, 6.65, 1.2, 6.25, 5.5, fill=RGBColor(0xFD,0xF8,0xFF), border=PURPLE, bw=1)
t = tb(sl, 6.8, 1.3, 5.95, 0.4)
p(t, "Model Selection Criteria (M2.5)", sz=17, bold=True, col=PURPLE)
t3 = tb(sl, 6.8, 1.75, 5.95, 4.8)
p(t3,  "Precond. gate  (CVD–HC direction)", sz=15, bold=True, col=BODY)
bl(t3, "Signed Cohen's d ≥ +0.5",     sz=14)
bl(t3, "Cells failing gate excluded",  sz=14)
p(t3,  "Primary: Parameter Stability", sz=15, bold=True, col=BODY, sb=8)
bl(t3, "Boundary rate < 50%",          sz=14)
bl(t3, "HC resample IQR (β_s,β_c) ↓", sz=14)
bl(t3, "σ-bin mode share ↑",           sz=14)
bl(t3, "Strict LOO range ↓",           sz=14)
p(t3,  "⚠ Specificity는 criterion 아님", sz=15, bold=True, col=AMBER, sb=8)
bl(t3, "HC FPR=100% → descriptive only",sz=14, col=AMBER)
p(t3,  "Fitting flow",                 sz=15, bold=True, col=BODY, sb=8)
bl(t3, "Grid search → z-score argmin", sz=14)
bl(t3, "RDM selects σ-bin;",           sz=14)
bl(t3, "γ/LOCO refine sub-bin",        sz=14)
bl(t3, "→ Pre-image (Brent, M4)",      sz=14)


# ── SLIDE 5: RQ1 R+C FAILS ────────────────────────────────────
sl = blank(prs)
hdr(sl, "RQ1 — R+C: Structurally Insufficient", 5,
    sub="Boundary saturation = model misspecification  (Wilson & Collins 2019)")

hdrs = ["Subj.", "Model", "Loss", "g / Δλ", "Boundary", "Verdict"]
rows = [
    ["S08 deutan", "R+C (JND_Lamb)",  "RDM_V1",        "Δλ=6.5nm, g=3.00", "100%  ⚠", "Failure"],
    ["S09 protan", "R+C (Boehm_low)","γ_all+RDM_V1",  "Δλ=3.0nm, g=2.95", "41%  ⚠",  "Near-sat."],
]
tbl = add_table(sl, hdrs, rows, 0.4, 1.25, 12.4, 1.35)
for j in range(6):
    for i in [1, 2]:
        c = tbl.cell(i, 5); c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0xFF,0xEE,0xEE)

box(sl, 0.4, 2.76, 12.3, 1.65, fill=LIGHT_R, border=REDF, bw=1)
tr = tb(sl, 0.6, 2.86, 12.0, 1.45)
p(tr, "Why is this a model-failure signal?", sz=15, bold=True, col=REDF)
bl(tr, "R+C: δθ=(2−g)·δθ_Machado  →  NO confusion-axis DOF (β_c)", sz=14)
bl(tr, "S08 g=3.0, boundary=100%: g stuck at grid ceiling, not an optimum", sz=14)
bl(tr, "S09 g=2.95, boundary=41%: near-saturation", sz=14)

box(sl, 0.4, 4.56, 12.3, 0.85, fill=RGBColor(0xEB,0xF7,0xEE), border=GREEN, bw=1)
ti = tb(sl, 0.6, 4.66, 12.0, 0.7)
p(ti, "결론: Rejection = structural DOF shortage  (NOT a literature g comparison)", sz=14, bold=True, col=GREEN)
bl(ti, "2-Component captures β_c — impossible with R+C's 1-DOF form", sz=14)

box(sl, 0.4, 5.56, 12.3, 0.95, fill=RGBColor(0xF9,0xF9,0xF9), border=DIVIDER, bw=1)
tn = tb(sl, 0.6, 5.65, 12.0, 0.8)
p(tn, "Note", sz=13, bold=True, col=GRAY)
bl(tn, "Lit. g (Boehm ~1.0–1.3, Tregillus ~1.1): paradigm differs → supporting evidence only", sz=13, col=GRAY)
bl(tn, "R+C kept in Supplement as retinal diagnostic decomposition", sz=13, col=GRAY)


# ── SLIDE 6: RQ2 FITTING RESULTS ──────────────────────────────
sl = blank(prs)
hdr(sl, "RQ2 — 2-Component Fitting Results & Stability", 6,
    sub="v6 PCA 45° RDM  |  N=300 HC resamples + strict 7-fold LOO")

hdrs2 = ["Label", "Subj.", "Loss", "(β_s,β_c)", "IQR", "Mode", "LOO β_c"]
rows2 = [
    ["βc-dom ★", "S08", "γ_OY + RDM_V2", "(+6,−42)", "(8,2)",  "~70%",  "[−46,−38]✓"],
    ["βc-rot ★", "S09", "γ_all+ RDM_V1", "(+2,+24)", "(0,0)",  "87.7%", "[24,24]✓"],
    ["βs-dom †", "S08", "γ_all + RDM_V1", "(+38,−10)","(12,4)","~50%",  "—"],
]
tbl2 = add_table(sl, hdrs2, rows2, 0.4, 1.25, 6.15, 1.4)
for j in range(len(hdrs2)):
    c = tbl2.cell(1, j); c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0xFF,0xF0,0xEE)
for j in range(len(hdrs2)):
    c = tbl2.cell(2, j); c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0xEE,0xF7,0xF3)

box(sl, 0.4, 2.8, 6.15, 2.0, fill=RGBColor(0xFF,0xF5,0xF4), border=S08C, bw=1)
ts08 = tb(sl, 0.55, 2.88, 5.85, 0.36)
p(ts08, "S08  βc-dom  — selection basis", sz=14, bold=True, col=S08C)
ts08b = tb(sl, 0.55, 3.28, 5.85, 1.45)
bl(ts08b, "IQR (8,2) < βs-dom (12,4): more stable", sz=14)
bl(ts08b, "LOO [−46,−38]: no zero-crossing",         sz=14)
bl(ts08b, "RDM ROI = V2 (geom. p=0.040), not post-hoc",sz=14)

box(sl, 0.4, 4.94, 6.15, 2.0, fill=RGBColor(0xF0,0xF9,0xF4), border=S09C, bw=1)
ts09 = tb(sl, 0.55, 5.02, 5.85, 0.36)
p(ts09, "S09  βc-rot  — interpretation", sz=14, bold=True, col=S09C)
ts09b = tb(sl, 0.55, 5.42, 5.85, 1.45)
bl(ts09b, "IQR=(0,0): deterministic  |  mode 87.7%", sz=14)
bl(ts09b, "β_c=+24° ← protan θ_conf=16°",            sz=14)
bl(ts09b, "⚠ PCA(2,+24) vs SRM(32,0): metric-dep.",  sz=14, col=AMBER)

img(sl, f"{V4COL}/p2_primary_4col_summary.png", l=6.75, t=1.2, w=6.35)


# ── SLIDE 7: RQ4+RQ5 NEURAL BENEFIT ───────────────────────────
sl = blank(prs)
hdr(sl, "RQ4+RQ5 — Neural Data Reveals What Behavior Cannot", 7,
    sub="Core rationale for fMRI-based filter  |  S09 β_c: behavior≈0, neural=+24°")

box(sl, 0.4, 1.2, 6.1, 1.85, fill=LBLUE, border=BLUE_H, bw=1)
th = tb(sl, 0.55, 1.28, 5.8, 0.35)
p(th, "Stability: behavior-only vs +neural (PCA)", sz=14, bold=True, col=BLUE_H)
hdrs_q4 = ["Subject", "Behav. IQR", "+Neural IQR", "Boundary Δ"]
rows_q4 = [
    ["S08", "(18,6)", "(8,2) ↓",  "23.0% → 9.3%"],
    ["S09", "(6,4)",  "(0,0) ↓",  "0% → 0%"],
]
add_table(sl, hdrs_q4, rows_q4, 0.4, 1.65, 6.1, 1.3)

box(sl, 6.85, 1.2, 6.1, 1.85, fill=RGBColor(0xFF,0xF8,0xE1), border=AMBER, bw=2)
tk = tb(sl, 7.0, 1.28, 5.85, 0.35)
p(tk, "★ Sub-09 critical finding", sz=14, bold=True, col=AMBER)
tk2 = tb(sl, 7.0, 1.65, 5.85, 1.3)
bl(tk2, "Behavior-only:  β_c ≈ 0  (invisible)", sz=15, bold=True)
bl(tk2, "+Neural:        β_c = +24° revealed",  sz=15, bold=True, col=S09C)
bl(tk2, "fMRI exposes cortical β_c absent in behavior",sz=14)

box(sl, 0.4, 3.2, 12.3, 1.55, fill=RGBColor(0xF5,0xF5,0xFF), border=BLUE_H, bw=1)
t5h = tb(sl, 0.55, 3.28, 12.0, 0.35)
p(t5h, "RQ5 — Behavioral vs. neural distortion direction", sz=15, bold=True, col=BLUE_H)
hdrs_q5 = ["Subject", "Behav. β_c", "Neural β_c", "Agreement", "Meaning"]
rows_q5 = [
    ["S08 deutan","NEG (7/8)","NEG (V2,V4)","✓ Agree",  "Triangulation"],
    ["S09 protan","≈ 0",      "POS (+24,V1)","✗ Disagree","Neural-only detection"],
]
add_table(sl, hdrs_q5, rows_q5, 0.4, 3.65, 12.3, 1.1)

box(sl, 0.4, 4.9, 12.3, 2.1, fill=RGBColor(0xEB,0xF7,0xEE), border=GREEN, bw=1)
timp = tb(sl, 0.6, 4.98, 12.0, 1.9)
p(timp, "fMRI 기반 필터의 정당성 (D2)", sz=15, bold=True, col=GREEN)
bl(timp, "Behavior-based filters (commercial): cannot capture S09's cortical β_c", sz=14)
bl(timp, "LOCO loss directly specifies corrective target — which stimulus direction", sz=14)
bl(timp, "  moves neural repr. toward HC-normal  (not just quantifying the gap)", sz=14, col=GRAY)
bl(timp, "→ Architecturally distinct from population-average retinal filter (Machado)", sz=14, bold=True, col=NAVY)


# ── SLIDE 8: RQ3 IDENTIFIABILITY ──────────────────────────────
sl = blank(prs)
hdr(sl, "RQ3 — Identifiability: Class ✓,  Precise Value ✗", 8,
    sub="Reliable resolution = sign quadrant only  |  Noise floor ~20°/25°")

box(sl, 0.4, 1.2, 5.5, 2.3, fill=RGBColor(0xEB,0xF7,0xEE), border=GREEN, bw=2)
tg = tb(sl, 0.55, 1.28, 5.2, 0.36)
p(tg, "Identifiable  ✓", sz=16, bold=True, col=GREEN)
tg2 = tb(sl, 0.55, 1.68, 5.2, 1.75)
bl(tg2, "Mechanism class (sign quadrant)",        sz=14, bold=True)
bl(tg2, "  S08: β_s>0, β_c<0  (deutan)",         sz=14)
bl(tg2, "  S09: β_s≈0, β_c>0  (protan)",         sz=14)
bl(tg2, "  All 3 metrics agree on quadrant",      sz=13, col=GRAY)
bl(tg2, "Dominant-axis direction",                sz=14, bold=True, sb=4)
bl(tg2, "  S08 β_c bias: 30.9° → 4.7° (fixed)",  sz=14)

box(sl, 0.4, 3.65, 5.5, 2.5, fill=RGBColor(0xFC,0xEC,0xEB), border=REDF, bw=2)
tr = tb(sl, 0.55, 3.73, 5.2, 0.36)
p(tr, "Non-identifiable  ✗", sz=16, bold=True, col=REDF)
tr2 = tb(sl, 0.55, 4.13, 5.2, 1.95)
bl(tr2, "Precise (β_s, β_c) absolute values",    sz=14, bold=True)
bl(tr2, "  Noise floor: β_s ~20°, β_c ~25°",     sz=14)
bl(tr2, "  f10° < 0.30  (3/3 FAIL)",              sz=14)
bl(tr2, "  Origin: 0/140 within 10°",             sz=14)
bl(tr2, "→ Descriptive embedding only",           sz=14, bold=True, col=REDF)

box(sl, 0.4, 6.28, 5.5, 0.65, fill=NAVY)
tn = tb(sl, 0.55, 6.35, 5.2, 0.55)
p(tn, "Both subjects: family-consistent quadrant across all metric variants", sz=13, bold=True, col=WHITE)

img(sl, f"{FIG_DIR}/fig_candidates_param_space.png", l=6.2, t=1.18, w=6.9)


# ── SLIDE 9: SPECIFICITY 0/3 ──────────────────────────────────
sl = blank(prs)
hdr(sl, "Specificity 검증: 0/3 FAIL — 정직한 한계", 9,
    sub="FDR-corrected 4-test (BH α=0.05)  |  Specificity ≠ selection criterion")

hdrs_sp = ["Test", "S08 ★", "S09 ★", "Verdict"]
rows_sp = [
    ["T1 param recovery", "f10°=0.26", "f10°=0.14", "FAIL ✗"],
    ["T2a origin (0,0)",  "0/140",    "0/140",      "FAIL ✗"],
    ["T2b HC pseudo-CVD", "NS",       "NS",          "FAIL ✗"],
    ["T2c label-perm",    "p=0.167",  "p=0.471",     "FAIL ✗"],
]
tbl_sp = add_table(sl, hdrs_sp, rows_sp, 0.4, 1.2, 5.5, 2.6)
for i in range(1, 5):
    c = tbl_sp.cell(i, 3); c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0xFF,0xEE,0xEE)

box(sl, 0.4, 3.95, 5.5, 1.45, fill=LIGHT_R, border=REDF, bw=1)
tb_lim = tb(sl, 0.55, 4.05, 5.2, 1.3)
p(tb_lim, "T2a (load-bearing):", sz=14, bold=True, col=REDF)
bl(tb_lim, "Zero-signal → argmin ~20–25° from origin", sz=14)
bl(tb_lim, "= pipeline built-in noise floor / bias",   sz=14)

box(sl, 0.4, 5.52, 5.5, 1.4, fill=RGBColor(0xEB,0xF7,0xEE), border=GREEN, bw=1)
tb_ok = tb(sl, 0.55, 5.62, 5.2, 1.2)
p(tb_ok, "주장 가능:", sz=14, bold=True, col=GREEN)
bl(tb_ok, "Real CVD minima 2.1–5.5× deeper (Exp17)", sz=14)
bl(tb_ok, "Mechanism class robust | pre-image 8/8 ✓",sz=14)

img(sl, f"{FIG_DIR}/fig_specificity_summary.png", l=6.2, t=1.15, w=6.95)

box(sl, 0.4, 7.25, 12.3, 0.19, fill=NAVY)
tks = tb(sl, 0.5, 7.27, 12.1, 0.17)
p(tks, "결론: all candidates = descriptive fit  |  Phase 3 = sole external validation", sz=12, bold=True, col=WHITE)


# ── SLIDE 10: PRE-IMAGE FILTER ────────────────────────────────
sl = blank(prs)
hdr(sl, "Pre-image 보정 필터 (R4)", 10,
    sub="Bijective 2-Component map → exact pre-image per hue  |  Brent refinement")

box(sl, 0.4, 1.2, 12.3, 1.0, fill=LBLUE, border=BLUE_H, bw=1)
tm = tb(sl, 0.6, 1.3, 12.0, 0.82)
p(tm, "Forward:  T(θ̃_k) = θ̃_k + δθ(θ̃_k) = θ_k    →    Filter:  δθ^filt_k = θ̃_k − θ_k",
  sz=16, bold=True, col=NAVY, align=PP_ALIGN.CENTER)
p(tm, "Criterion: residual < 0.001° for all 8 hues  —  failure rejects subject–model pairing",
  sz=13, col=GRAY, align=PP_ALIGN.CENTER)

hdrs_fi = ["Subject", "Candidate", "(β_s, β_c)", "8/8?", "Note"]
rows_fi = [
    ["S08 (deutan)", "βc-dom ★",  "(+6, −42)",  "✓", "Primary filter"],
    ["S08 (deutan)", "βs-dom †",  "(+38, −10)", "✓", "Supplement (parallel)"],
    ["S09 (protan)", "βc-rot ★",  "(+2, +24)",  "✓", "Primary — prev. draft −22° corrected"],
]
add_table(sl, hdrs_fi, rows_fi, 0.4, 2.35, 12.3, 1.55)

box(sl, 0.4, 4.06, 12.3, 0.8, fill=RGBColor(0xFF,0xF8,0xE1), border=AMBER, bw=1)
tb_bd = tb(sl, 0.6, 4.14, 12.0, 0.65)
p(tb_bd, "Model-class bound:", sz=14, bold=True, col=AMBER)
bl(tb_bd, "|δθ(45°)| ≤ |0.71·β̂_s − 0.26·β̂_c|  →  distortions > ~51° outside current model class",
  sz=14)

box(sl, 0.4, 5.02, 12.3, 1.45, fill=LBLUE, border=BLUE_H, bw=1)
tb_r = tb(sl, 0.6, 5.12, 12.0, 1.25)
p(tb_r, "남은 실행 항목", sz=14, bold=True, col=BLUE_H)
bl(tb_r, "E1: Recompute δθ 8-vec via two_comp.py:forward_2comp (S09 direction fix)", sz=14)
bl(tb_r, "R5: 2AFC behavioral + 2nd-session fMRI (Phase 3, TBD)", sz=14)

box(sl, 0.4, 6.6, 12.3, 0.4, fill=RGBColor(0xEC,0xF0,0xF1))
tp = tb(sl, 0.6, 6.65, 12.0, 0.3)
p(tp, "Pipeline:  B (fitting)  →  ★ C (pre-image)  →  [pending] D (Phase 3 validation)", sz=14, col=NAVY)


# ── SLIDE 11: DISCUSSION ──────────────────────────────────────
sl = blank(prs)
hdr(sl, "Discussion — fMRI 기반 필터 정당성 + 공통 피질 모델", 11,
    sub="D2: neural advantage  ·  D4: common model  ·  D3: etiology  ·  D6: upstream alt.")

box(sl, 0.4, 1.2, 8.1, 2.85, fill=RGBColor(0xF0,0xF6,0xFF), border=BLUE_H, bw=2)
td2 = tb(sl, 0.55, 1.28, 7.8, 0.36)
p(td2, "D2 — Neural reveals what behavior cannot", sz=16, bold=True, col=BLUE_H)
td2b = tb(sl, 0.55, 1.68, 7.8, 2.28)
bl(td2b, "S09: behav-only β_c≈0 → +neural β_c=+24° (cortical rotation exposed)",
  sz=14, bold=True, col=NAVY)
bl(td2b, "Commercial filters (behavior-based): miss cortical β_c",   sz=14)
bl(td2b, "LOCO loss specifies corrective target — prescribes stimulus",sz=14)
bl(td2b, "  shift direction to align neural repr. with HC-normal",    sz=14, col=GRAY)
bl(td2b, "→ Architecturally distinct from population-avg retinal filt.",sz=14, bold=True, col=NAVY)

box(sl, 0.4, 4.18, 8.1, 2.82, fill=RGBColor(0xF4,0xF0,0xFF), border=PURPLE, bw=2)
td4 = tb(sl, 0.55, 4.26, 7.8, 0.36)
p(td4, "D4 — 공통 피질 모델, 개인화 파라미터", sz=16, bold=True, col=PURPLE)
td4b = tb(sl, 0.55, 4.66, 7.8, 2.26)
bl(td4b, "One 2-Component model spans both CVD subtypes",       sz=14, bold=True, col=NAVY)
bl(td4b, "Difference is (β_s,β_c), NOT mechanism class:",       sz=14)
bl(td4b, "  Deutan: large confusion rotation  (β_c=−42°)",     sz=14, col=S08C)
bl(td4b, "  Protan: small confusion rotation  (β_c=+24°)",     sz=14, col=S09C)
bl(td4b, "Stronger than two-mechanism explanation",             sz=14, bold=True)

box(sl, 8.85, 1.2, 4.05, 2.85, fill=RGBColor(0xFF,0xF9,0xEC), border=AMBER, bw=1)
td3 = tb(sl, 9.0, 1.28, 3.75, 0.36)
p(td3, "D3 — Etiology (brief)", sz=15, bold=True, col=AMBER)
td3b = tb(sl, 9.0, 1.68, 3.75, 2.28)
bl(td3b, "R+C framework motivates forward modelling", sz=14)
bl(td3b, "Both subjects: R+C structural saturation", sz=14)
bl(td3b, "→ cortical distortion locus",              sz=14, bold=True)
bl(td3b, "Both need cortical opponent-rotation",     sz=14, sb=4)
bl(td3b, "(Emery 2021/23 convention; no '21.4°')",   sz=13, col=GRAY)
bl(td3b, "⚠ Structural detail in R1, not here",      sz=13, col=GRAY)

box(sl, 8.85, 4.18, 4.05, 2.82, fill=RGBColor(0xF5,0xF5,0xF5), border=DIVIDER, bw=1)
td6 = tb(sl, 9.0, 4.26, 3.75, 0.36)
p(td6, "D6 — Upstream-input (retain)", sz=14, bold=True, col=GRAY)
td6b = tb(sl, 9.0, 4.66, 3.75, 2.2)
bl(td6b, "Is hV4 damage passive V1–V3 cascade?", sz=14)
bl(td6b, "LORO preservation refutes this:",       sz=14)
bl(td6b, "  discrimination intact → V1–V3 OK",   sz=14, col=GRAY)
bl(td6b, "  hV4 selectively loses interpolation", sz=14, col=GRAY)
bl(td6b, "Keep for reviewer preemption",          sz=13, bold=True, col=GRAY, sb=4)


# ── SLIDE 12: LIMITATIONS + NEXT STEPS ───────────────────────
sl = blank(prs)
hdr(sl, "한계 및 다음 단계 (D7 + Phase 3)", 12,
    sub="Three-theme limitation structure  |  Phase 3 = sole verification path")

box(sl, 0.35, 1.2, 4.0, 4.5, fill=LIGHT_R, border=REDF, bw=1)
ta_h = tb(sl, 0.5, 1.28, 3.7, 0.36)
p(ta_h, "Theme A — Identifiability", sz=15, bold=True, col=REDF)
ta_b = tb(sl, 0.5, 1.68, 3.7, 3.95)
bl(ta_b, "Class (quadrant): robust ✓",  sz=14, col=GREEN, bold=True)
bl(ta_b, "Precise value: non-ident. ✗", sz=14, col=REDF, bold=True)
bl(ta_b, "4-test summary:",             sz=14, bold=True, sb=5)
bl(ta_b, "  T1 recovery: 0/3 FAIL",    sz=13, col=REDF)
bl(ta_b, "  T2a origin:  0/140 hit",   sz=13, col=REDF)
bl(ta_b, "  T2b/2c: 0/3 pass",         sz=13, col=REDF)
bl(ta_b, "Noise floor: β_s~20°,β_c~25°",sz=14, sb=5)
bl(ta_b, "S09 metric-dependence:",      sz=14, sb=5)
bl(ta_b, "  PCA(2,+24) vs SRM(32,0)",  sz=13, col=AMBER)
bl(ta_b, "  → mechanism class changes",sz=13, col=AMBER)

box(sl, 4.65, 1.2, 4.0, 4.5, fill=LIGHT_G, border=GREEN, bw=1)
tb_h = tb(sl, 4.8, 1.28, 3.7, 0.36)
p(tb_h, "Theme B — Sample / OOS", sz=15, bold=True, col=GREEN)
tb_b = tb(sl, 4.8, 1.68, 3.7, 3.95)
bl(tb_b, "CVD N=2: no CVD LOO",         sz=14)
bl(tb_b, "All OOS axes = HC pool only", sz=14)
bl(tb_b, "CVD JND: N=1 per pair",       sz=14)
bl(tb_b, "Held-out focal pair:",         sz=14, sb=5)
bl(tb_b, "  same CVD obs. re-used",     sz=14, col=AMBER)
bl(tb_b, "  HC norm. changes only",     sz=14, col=AMBER)
bl(tb_b, "→ Generalization = Phase 3", sz=14, bold=True, col=NAVY, sb=5)

box(sl, 8.95, 1.2, 4.0, 4.5, fill=RGBColor(0xFF,0xF8,0xE1), border=AMBER, bw=1)
tc_h = tb(sl, 9.1, 1.28, 3.7, 0.36)
p(tc_h, "Theme C — Model Framework", sz=15, bold=True, col=AMBER)
tc_b = tb(sl, 9.1, 1.68, 3.7, 3.95)
bl(tc_b, "Z-score composite:", sz=14)
bl(tc_b, "  1-pair γ = 8-pair γ weight",sz=14, col=GRAY)
bl(tc_b, "  (intended normalization)",   sz=13, col=GRAY)
bl(tc_b, "R+C 1-DOF:",                  sz=14, sb=5)
bl(tc_b, "  structural constraint,",    sz=14, col=GRAY)
bl(tc_b, "  not lit. comparison",       sz=13, col=GRAY)
bl(tc_b, "PCA-RDM adoption:",           sz=14, sb=5)
bl(tc_b, "  Cycle 5: 2× separation",   sz=14, col=GRAY)
bl(tc_b, "  S09 stability>SRM",        sz=14, col=GRAY)
bl(tc_b, "  (less established metric)",sz=13, col=AMBER)

box(sl, 0.35, 5.85, 12.6, 1.45, fill=LBLUE, border=BLUE_H, bw=2)
t3_h = tb(sl, 0.55, 5.93, 12.2, 0.36)
p(t3_h, "Phase 3 — Sole Verification Path", sz=16, bold=True, col=BLUE_H)
t3_b = tb(sl, 0.55, 6.33, 12.2, 0.9)
bl(t3_b, "2AFC behavioral: filter vs. no-filter → JND reduction (only cortical 2-comp filter predicted to pass)",
  sz=14)
bl(t3_b, "2nd-session fMRI: SRM/LOCO validation  |  Baseline = Experiment 1",
  sz=14)

# ── SAVE ──────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved: {OUT}  ({len(prs.slides)} slides)")
